import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    assets_dir = r"c:\Users\shaur\OneDrive\web2\skillpath\ppt_assets"

    # Color Palette
    COLOR_TITLE = RGBColor(15, 23, 42)      # #0F172A Slate 900
    COLOR_BODY = RGBColor(30, 41, 59)       # #1E293B Slate 800
    COLOR_MUTED = RGBColor(71, 85, 105)     # #475569 Slate 600
    COLOR_PRIMARY = RGBColor(37, 99, 235)   # #2563EB Blue 600
    COLOR_ACCENT = RGBColor(234, 88, 12)    # #EA580C Orange 600
    COLOR_EMERALD = RGBColor(16, 185, 129)  # #10B981 Emerald 500
    COLOR_BG_CARD = RGBColor(248, 250, 252) # #F8FAFC Slate 50
    COLOR_BORDER = RGBColor(203, 213, 225)  # #CBD5E1 Slate 300
    COLOR_WHITE = RGBColor(255, 255, 255)

    def add_slide_bg(slide, slide_num):
        if slide_num in range(2, 10):
            bg_path = os.path.join(assets_dir, f"slide_{slide_num}_clean.png")
        else:
            bg_path = os.path.join(assets_dir, f"slide_{slide_num}.png")
            
        if os.path.exists(bg_path):
            slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    def add_card(slide, left, top, width, height, bg_color=COLOR_BG_CARD, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ----------------------------------------------------
    # SLIDE 1: Title Slide Overlay
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide1, 1)

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.333), Inches(2.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "SkillPath — Autonomous Inclusive Workforce Orchestrator"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Closing the $5.5T Skills Gap with Skills-First Intelligence & SAP Talent Intelligence Hub Alignment"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(224, 231, 255)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 2: Problem Statement & Team Details
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide2, 2)

    # Team Name & Problem Statement Header Boxes
    add_card(slide2, Inches(1.1), Inches(1.7), Inches(5.4), Inches(1.0), bg_color=RGBColor(238, 242, 255), border_color=COLOR_PRIMARY)
    tb_tn = slide2.shapes.add_textbox(Inches(1.2), Inches(1.75), Inches(5.2), Inches(0.9))
    tf_tn = tb_tn.text_frame
    tf_tn.word_wrap = True
    p = tf_tn.paragraphs[0]
    p.text = "TEAM NAME (AS SUBMITTED ON PORTAL)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED
    p2 = tf_tn.add_paragraph()
    p2.text = "Team SkillPath"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_PRIMARY
    p2.space_before = Pt(4)

    add_card(slide2, Inches(6.7), Inches(1.7), Inches(5.533), Inches(1.0), bg_color=RGBColor(238, 242, 255), border_color=COLOR_PRIMARY)
    tb_ps = slide2.shapes.add_textbox(Inches(6.8), Inches(1.75), Inches(5.333), Inches(0.9))
    tf_ps = tb_ps.text_frame
    tf_ps.word_wrap = True
    p = tf_ps.paragraphs[0]
    p.text = "HACKATHON THEME & PROBLEM STATEMENT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED
    p2 = tf_ps.add_paragraph()
    p2.text = "Inclusive Workforce: AI Multi-Agent Career Orchestrator & Fair Reskilling"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_before = Pt(4)

    # Team Members Table
    table_shape = slide2.shapes.add_table(5, 4, Inches(1.1), Inches(2.85), Inches(11.133), Inches(3.9))
    table = table_shape.table

    table.columns[0].width = Inches(0.9)   # S.No
    table.columns[1].width = Inches(2.7)   # NAME
    table.columns[2].width = Inches(1.5)   # GENDER
    table.columns[3].width = Inches(6.033) # AREA OF EXPERTISE

    headers = ["S.No", "NAME", "GENDER", "AREA OF EXPERTISE"]
    for col_idx, h_text in enumerate(headers):
        table.cell(0, col_idx).text = h_text

    members = [
        ("1", "Shaurya", "Male", "AI Multi-Agent Architect, DAG Swarm & SAP Integration"),
        ("2", "Harshit", "Male", "Frontend Engineering, Inclusive UI/UX & Cloud Infrastructure"),
        ("3", "", "", ""),
        ("4", "", "", "")
    ]

    for row_idx, m in enumerate(members):
        r = row_idx + 1
        table.cell(r, 0).text = m[0]
        table.cell(r, 1).text = m[1]
        table.cell(r, 2).text = m[2]
        table.cell(r, 3).text = m[3]

    for r_i, row in enumerate(table.rows):
        for c_i, cell in enumerate(row.cells):
            cell.fill.solid()
            if r_i == 0:
                cell.fill.fore_color.rgb = COLOR_PRIMARY
            else:
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if r_i % 2 == 1 else COLOR_WHITE

            for p in cell.text_frame.paragraphs:
                p.font.name = "Arial"
                if r_i == 0:
                    p.font.size = Pt(13)
                    p.font.bold = True
                    p.font.color.rgb = COLOR_WHITE
                else:
                    p.font.size = Pt(12)
                    p.font.color.rgb = COLOR_BODY

    # ----------------------------------------------------
    # SLIDE 3: Title Name & Proposed Solution
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide3, 3)

    tb_h3 = slide3.shapes.add_textbox(Inches(1.1), Inches(1.0), Inches(11.133), Inches(0.9))
    tf_h3 = tb_h3.text_frame
    p3_1 = tf_h3.paragraphs[0]
    p3_1.text = "SkillPath — The Inclusive Workforce Orchestrator"
    p3_1.font.size = Pt(24)
    p3_1.font.bold = True
    p3_1.font.color.rgb = COLOR_TITLE

    p3_2 = tf_h3.add_paragraph()
    p3_2.text = "Proposed Solution Architecture:"
    p3_2.font.size = Pt(16)
    p3_2.font.bold = True
    p3_2.font.color.rgb = COLOR_PRIMARY
    p3_2.space_before = Pt(4)

    card_w = Inches(3.5)
    card_h = Inches(4.7)
    card_top = Inches(2.05)

    cards_s3 = [
        ("🌟 Gap Alchemy & Skills Discovery", [
            "• Translates informal lived experience (caregiving, operations, self-study) into high-value tech skills.",
            "• Enforces a 0% Career Break Penalty policy, surfacing latent ability over institution pedigree.",
            "• Connects adjacent capabilities through a multi-dimensional skills knowledge graph."
        ]),
        ("🤖 6-Agent Autonomous Swarm", [
            "• Skills Discovery, Market Intelligence, Learning Pathway, Inclusive Matcher, Employer Court, Bias Auditor.",
            "• Stepped Bridge Role Ladders: builds viable stepping-stones rather than unrealistic career leaps.",
            "• Real-time doubt resolution with Human-in-the-Loop decision agency."
        ]),
        ("🏢 SAP SuccessFactors & Joule Bridge", [
            "• Standard Skills Portfolio JSON export matching SAP Talent Intelligence Hub taxonomy.",
            "• Pre-packaged handoff payloads for SAP Joule Career, HR Service, and People Intelligence Agents.",
            "• Certified 10-Dimension Fairness & Bias Audit Ledger with verifiable A+ grade."
        ])
    ]

    for idx, (title, bullets) in enumerate(cards_s3):
        left_pos = Inches(1.1 + idx * 3.8)
        add_card(slide3, left_pos, card_top, card_w, card_h)

        tb = slide3.shapes.add_textbox(left_pos + Inches(0.15), card_top + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TITLE

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = b
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = COLOR_BODY
            p_b.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 4: Outline of Unique & Innovative Solution
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide4, 4)

    s4_innovations = [
        ("🛡️ 10-Dimension Bias & Gap Immunity Engine", "Evaluates recommendations against 10 fairness dimensions including PwD accessibility, tier-2/3 access, caregiving break protection, and degree-blind skills verification."),
        ("🗺️ Stepped Bridge-Role Ladders (Dijkstra Graph)", "Rather than demanding an impossible 6-month jump from back-office data entry to Senior AI, Pathfinder builds stepped, paid bridge roles (Data Associate ➔ IT Support ➔ Cloud QA)."),
        ("⚖️ Employer Court / JD Accessibility Audit", "Scans employer job postings for exclusionary jargon ('rockstar', 'ivy league mandate', 'unbroken tenure') and generates actionable HR accommodation recommendations."),
        ("⚡ DAG Wave Parallel Execution (< 2.5s)", "Executes 14 specialized sub-agents in 4 parallel waves using SHA-256 caching and fallback resilience, delivering enterprise speed with deep agentic reasoning.")
    ]

    grid_w = Inches(5.3)
    grid_h = Inches(2.25)
    
    positions_s4 = [
        (Inches(1.1), Inches(1.85)),
        (Inches(6.9), Inches(1.85)),
        (Inches(1.1), Inches(4.35)),
        (Inches(6.9), Inches(4.35)),
    ]

    for idx, (title, desc) in enumerate(s4_innovations):
        left_pos, top_pos = positions_s4[idx]
        add_card(slide4, left_pos, top_pos, grid_w, grid_h)

        tb = slide4.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), grid_w - Inches(0.4), grid_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_BODY
        p_d.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 5: Technical Approach & Swarm Architecture
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide5, 5)

    tech_pillars = [
        ("🖥️ Frontend & UI Layer", [
            "Next.js 15 App Router & React 19",
            "TypeScript 5 Strict (End-to-End)",
            "Tailwind CSS + Glassmorphic UI",
            "Interactive Copilot & Dynamic Re-run",
            "4 One-Click Live Persona Presets"
        ]),
        ("🤖 6-Agent Swarm Core", [
            "DAG Parallel Wave Scheduler (<2.5s)",
            "Gap-Skill Translation Sub-Engine",
            "Dijkstra Shortest Bridge Pathfinder",
            "10-Dimension Inclusion Rules Engine",
            "Employer Readiness Jargon Auditor"
        ]),
        ("🏢 Enterprise & SAP Bridge", [
            "SAP Talent Intelligence Hub JSON Export",
            "SAP Joule Agent Handoff Simulation",
            "Certified Bias Audit Certificate Export",
            "Key Pool Rotation & Model Fallback",
            "Zero-Trust Serverless Security"
        ]),
        ("📊 Datasets & Intelligence", [
            "5,000+ JD MVC Frequency Models",
            "Indian Tier-1/2/3 Regional Hub Matrix",
            "Live Skill Freshness & Decay Vectors",
            "Curated Returnship Programs DB",
            "WCAG 2.1 PwD Accessibility Mappings"
        ])
    ]

    col_w = Inches(2.65)
    col_h = Inches(4.8)

    for idx, (title, items) in enumerate(tech_pillars):
        left_pos = Inches(1.1 + idx * 2.85)
        add_card(slide5, left_pos, Inches(1.85), col_w, col_h)

        tb = slide5.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.95), col_w - Inches(0.3), col_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TITLE

        for item in items:
            p_i = tf.add_paragraph()
            p_i.text = f"• {item}"
            p_i.font.size = Pt(11.5)
            p_i.font.color.rgb = COLOR_BODY
            p_i.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 6: Feasibility, Viability & SAP Alignment
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide6, 6)

    viability_cards = [
        ("🏢 SAP SuccessFactors Synergy", "Native architectural alignment with SAP Talent Intelligence Hub and Joule Agents (Career & Talent Development, HR Service, People Intelligence) enables seamless enterprise talent deployment."),
        ("⚡ Production Feasibility", "Fully functional Next.js 15 application with sub-2.5s multi-agent execution, local deterministic reflex layer, and SHA-256 in-memory caching."),
        ("💰 Cost & Latency Optimization", "Hybrid deterministic + LLM wave architecture reduces API operational costs by >65% compared to monolithic LLM chains while preventing rate limit failures."),
        ("🛡️ Governance & Ethical AI", "Certified 10-dimension Bias Audit Agent enforces verifiable compliance with equal-opportunity standards and zero career break penalties.")
    ]

    for idx, (title, desc) in enumerate(viability_cards):
        left_pos, top_pos = positions_s4[idx]
        add_card(slide6, left_pos, top_pos, grid_w, grid_h)

        tb = slide6.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), grid_w - Inches(0.4), grid_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_BODY
        p_d.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 7: Social Dimension & Persona Impact
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide7, 7)

    impact_groups = [
        ("🌟 Women Returners & Gaps", [
            "• Caregiving and family breaks protected with 0% penalty.",
            "• Lived operations & budgeting translated into IT project skills.",
            "• Automatic matching with Microsoft LEAP, IBM & TCS returnships."
        ]),
        ("🚀 Tier-2/3 & First-Gen Talent", [
            "• Demystifies the hidden curriculum of career navigation.",
            "• Regional scouting for candidates in Lucknow, Jaipur, Gorakhpur.",
            "• Replaces institution prestige filters with verified skills."
        ]),
        ("♿ Displaced Workers & PwD", [
            "• Stepped 8-week bridge ladders for automated routine workers.",
            "• Factor screen-reader & async remote accommodations.",
            "• Salary ROI modeling for realistic reskilling upside."
        ])
    ]

    for idx, (title, bullets) in enumerate(impact_groups):
        left_pos = Inches(1.1 + idx * 3.8)
        add_card(slide7, left_pos, card_top, card_w, card_h)

        tb = slide7.shapes.add_textbox(left_pos + Inches(0.15), card_top + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TITLE

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = b
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = COLOR_BODY
            p_b.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 8: Business & Enterprise Ecosystem Model
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide8, 8)

    biz_models = [
        ("👤 Free B2C Candidate Lifeline", "Lifelong free access for job seekers, returners, and first-gen students to discover skills, get gap protection, and access curated free roadmaps."),
        ("🏢 Enterprise SAP Ecosystem SaaS", "Integration subscription for global enterprises (Grundfos, Capgemini, IBM) deploying SAP SuccessFactors to ingest verified inclusive candidate portfolios."),
        ("🏛️ University & State Placement SaaS", "Cohort analytics and batch readiness monitoring for tier-2/3 universities and state skill development missions (PMKVY, NSDC)."),
        ("⚖️ Employer Bias Audit & HR API", "Automated compliance API auditing corporate job postings for exclusionary filters, promoting ESG and DEI hiring benchmarks.")
    ]

    for idx, (title, desc) in enumerate(biz_models):
        left_pos, top_pos = positions_s4[idx]
        add_card(slide8, left_pos, top_pos, grid_w, grid_h)

        tb = slide8.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), grid_w - Inches(0.4), grid_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_BODY
        p_d.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 9: Macro Data, Research & Citations
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide9, 9)

    ref_cards = [
        ("📊 IDC Enterprise Research 2026", "Over 90% of global enterprises face critical skills shortages, creating a $5.5 trillion loss risk rooted in knowledge mismatches rather than labor shortages."),
        ("🇮🇳 India Skills Report 2026 & TCS", "Employability reached 56.35% with tier-2/3 talent emergence, juxtaposed against 12,000 TCS layoffs citing AI-driven skill mismatches."),
        ("💼 PwC & Anthropic Impact Reports", "Workers with AI fluency command a 56% wage premium, while non-AI workers risk medium-term obsolescence without accessible reskilling."),
        ("🌐 World Economic Forum (WEF)", "170M new roles vs 92M displaced (78M net gain), but 59% of the workforce requires urgent, targeted retraining to prevent displacement.")
    ]

    for idx, (title, desc) in enumerate(ref_cards):
        left_pos, top_pos = positions_s4[idx]
        add_card(slide9, left_pos, top_pos, grid_w, grid_h)

        tb = slide9.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), grid_w - Inches(0.4), grid_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_BODY
        p_d.space_before = Pt(6)

    # Save presentation
    output_path = r"c:\Users\shaur\OneDrive\web2\skillpath\SkillPath_byteBuilt_PitchDeck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
